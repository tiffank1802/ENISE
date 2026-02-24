from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

# --- COLORS & STYLES ---
ENISE_BLUE = RGBColor(0, 51, 102)    # #003366
ENISE_LIGHT = RGBColor(0, 86, 179)   # #0056b3
GOLD_ACCENT = RGBColor(255, 215, 0)  # #FFD700
GREY_TEXT = RGBColor(51, 65, 85)     # #334155
BG_GREY = RGBColor(241, 245, 249)    # #F1F5F9

def apply_title_style(shape, text):
    shape.text = text
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.font.name = 'Arial'
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.color.rgb = ENISE_BLUE
    paragraph.alignment = PP_ALIGN.LEFT

def apply_body_style(text_frame):
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = GREY_TEXT
        paragraph.space_after = Pt(14)

# --- SLIDE 1: TITLE ---
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)

# Background decoration (Manual implementation of the theme)
left = top = Inches(0)
width = Inches(10)
height = Inches(0.15)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(7.35), width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = ENISE_BLUE
shape.line.fill.background()

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "PhD Interview: MAD-SIM Project"
subtitle.text = "Kevin TONGUE\nModelling powder flows in agitated devices\n\nIMT Mines Albi | LUT University | École des Mines de Saint-Étienne"

# Styling Title Slide
title.text_frame.paragraphs[0].font.name = "Arial"
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = ENISE_BLUE

subtitle.text_frame.paragraphs[0].font.name = "Arial"
subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.color.rgb = GREY_TEXT

# Speaker Notes
slide.notes_slide.notes_text_frame.text = (
    "Opening: Good morning. I am Kevin Tongue, a 5th-year student at Centrale Lyon-ENISE. "
    "I am honored to present my candidacy for the MAD-SIM project.\n\n"
    "Context: Mention the consortium immediately: Albi (experimental), St-Étienne (materials), LUT (numerical)."
)


# --- SLIDE 2: AGENDA ---
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "Presentation Agenda")

content = slide.placeholders[1]
content.text = "Academic Background: Training & expertise in numerical mechanics.\n" \
               "Research Experience: From fatigue analysis to granular segregation.\n" \
               "Understanding the Topic: The challenge of powder rheology.\n" \
               "Vision 2026-2029: A 3-year roadmap.\n" \
               "Profile Match: Skills & Motivations."

apply_body_style(content.text_frame)


# --- SLIDE 3: ACADEMIC BACKGROUND ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "Academic Background")

content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear() 

p = text_frame.add_paragraph()
p.text = "Centrale Lyon - ENISE"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ENISE_BLUE

p = text_frame.add_paragraph()
p.text = "Engineering Degree: Mechanical Engineering (Design & Simulation)."
p.level = 0
p.font.size = Pt(20)

p = text_frame.add_paragraph()
p.text = "Master 2 Research: Numerical Solid Mechanics."
p.level = 0
p.font.size = Pt(20)

p = text_frame.add_paragraph()
p.text = "Developed expertise in multi-scale modeling, multiphysics coupling, and digital twins."
p.level = 1
p.font.italic = True


# --- SLIDE 4: RESEARCH EXPERIENCE ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "Research Experience")

content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "Fatigue Analysis (Current)"
p.font.bold = True
p.font.color.rgb = ENISE_BLUE

p = text_frame.add_paragraph()
p.text = "Optimization of the Dang Van criterion for industrial structures under complex loading."
p.font.size = Pt(18)

p = text_frame.add_paragraph()
p.text = "Skill: Validating numerical models against real constraints."
p.level = 1
p.font.italic = True

p = text_frame.add_paragraph()
p.text = "" # Spacer

p = text_frame.add_paragraph()
p.text = "The 'Stepping Stone' (March 2026)"
p.font.bold = True
p.font.color.rgb = ENISE_BLUE

p = text_frame.add_paragraph()
p.text = "Internship: Granular Segregation via Markov Chains."
p.font.size = Pt(18)

p = text_frame.add_paragraph()
p.text = "Skill: Starting the PhD with a solid foundation in DEM and segregation physics."
p.level = 1
p.font.italic = True

# Notes
slide.notes_slide.notes_text_frame.text = (
    "Highlight the internship as 'Phase 0'. You are arriving with the tools already in hand. "
    "Mention Dang Van to show you understand complex validation."
)


# --- SLIDE 5: UNDERSTANDING MAD-SIM ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "Understanding the MAD-SIM Challenge")

content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "Systemic Powder Rheology"
p.font.bold = True
p.font.color.rgb = ENISE_BLUE

p = text_frame.add_paragraph()
p.text = "Powders oscillate between solid, fluid, and gas states. Industrial sizing is currently too empirical."
p.font.size = Pt(20)

p = text_frame.add_paragraph()
p.text = "Goal: Rationalize agitation operations."
p.level = 1

p = text_frame.add_paragraph()
p.text = "Method: Use pilot as a systemic rheometer."
p.level = 1

p = text_frame.add_paragraph()
p.text = "Result: Identify invariant law μ(I) linking micro-collisions to macro-flow."
p.level = 1

# Placeholder for Image
left = Inches(7)
top = Inches(2)
width = Inches(2.5)
height = Inches(2.5)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.text = "[Insert Image from Thesis Subject Page 2 here]"


# --- SLIDE 6: VISION (TIMELINE) ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "Vision for the Three-Year Thesis")

content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "Year 1: Albi (Breaking the 'Black Box')"
p.font.bold = True
p.font.color.rgb = ENISE_BLUE
p = text_frame.add_paragraph()
p.text = "Physical observations on 9L pilot & initial DEM simulations."
p.level = 1

p = text_frame.add_paragraph()
p.text = "Year 2: Finland (International Immersion)"
p.font.bold = True
p.font.color.rgb = ENISE_BLUE
p = text_frame.add_paragraph()
p.text = "Advanced CFD & Numerical Methods at LUT University (Prof. Hämäläinen)."
p.level = 1

p = text_frame.add_paragraph()
p.text = "Year 3: Validation (Actionability)"
p.font.bold = True
p.font.color.rgb = ENISE_BLUE
p = text_frame.add_paragraph()
p.text = "Experimental validation & Industrial scale-up."
p.level = 1


# --- SLIDE 7: INTERNATIONAL EXCELLENCE ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "International Excellence")

content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.text = "Collaboration with LUT University (Finland)"
text_frame.paragraphs[0].font.bold = True
text_frame.paragraphs[0].font.color.rgb = ENISE_BLUE

p = text_frame.add_paragraph()
p.text = "The project leverages the expertise of the School of Engineering Sciences."
p.level = 0

p = text_frame.add_paragraph()
p.text = "Under Prof. Jari Hämäläinen, I will master cutting-edge CFD methods to bridge the gap between granular physics and industrial simulation."
p.level = 0

# Placeholder for LUT Image
left = Inches(6)
top = Inches(3)
width = Inches(3.5)
height = Inches(2.5)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.text = "[Insert LUT Campus Image]"


# --- SLIDE 8: SKILLS ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "A Profile at the Interface")

# Create Table
rows = 5
cols = 3
left = Inches(0.5)
top = Inches(2.0)
width = Inches(9.0)
height = Inches(3.0)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)
table.columns[2].width = Inches(3.0)

# Headers
headers = ["Domain", "Key Skills", "Tools"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = ENISE_BLUE
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    cell.text_frame.paragraphs[0].font.bold = True

# Data
data = [
    ["Modeling", "Solid Mechanics, Granular Media", "DEM, CFD, FEM"],
    ["Numerical", "Scientific Computing, Algorithms", "Python, C++, Fortran"],
    ["Simulation", "Data Analysis, Digital Twins", "Matlab (Certified), Simulink"],
    ["Experimental", "Industrial Pragmatism", "Fablab, Maintenance"]
]

for row_idx, row_data in enumerate(data):
    for col_idx, cell_value in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_value
        cell.text_frame.paragraphs[0].font.size = Pt(16)


# --- SLIDE 9: MOTIVATION ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "Actionability through Curiosity")

content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "Goal: Develop Digital Twins for the pharmaceutical industry."
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = ENISE_BLUE
p.alignment = PP_ALIGN.CENTER

p = text_frame.add_paragraph()
p.text = ""

p = text_frame.add_paragraph()
p.text = '"Walker, there is no path. The path is made by walking."'
p.font.size = Pt(32)
p.font.italic = True
p.font.color.rgb = GREY_TEXT
p.alignment = PP_ALIGN.CENTER

p = text_frame.add_paragraph()
p.text = "— Herbert Simon"
p.font.size = Pt(24)
p.font.color.rgb = GOLD_ACCENT
p.alignment = PP_ALIGN.RIGHT

# Notes
slide.notes_slide.notes_text_frame.text = (
    "This is your Strong Value. Connect the 'Land of Curious' (LUT motto) with your desire to build actionable industrial tools."
)


# --- SLIDE 10: Q&A ---
slide = prs.slides.add_slide(slide_layout)
apply_title_style(slide.shapes.title, "Questions?")

content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "Thank you for your time and attention."
p.font.size = Pt(24)
p.alignment = PP_ALIGN.CENTER

p = text_frame.add_paragraph()
p.text = "Kevin TONGUE"
p.font.bold = True
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

p = text_frame.add_paragraph()
p.text = "tonguekevin00@gmail.com"
p.font.size = Pt(18)
p.alignment = PP_ALIGN.CENTER


# Save
prs.save('kevin_tongue_madsim_presentation.pptx')